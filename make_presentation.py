"""
Generate CEO-facing PowerPoint for System Dynamics Model Factory (AeroDyn).
Run: python make_presentation.py
Output: AeroDyn_System_Dynamics_Model_Factory.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Slide dimensions (default 10x7.5 in)
SLIDE_W = 10.0
SLIDE_H = 7.5

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left, top, w, h = Inches(0.5), Inches(2.2), Inches(9), Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, w, h)
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x47, 0x5d)
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1.5))
        tx2.text_frame.word_wrap = True
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x4a, 0x4a, 0x4a)
    return slide

def add_section_slide(prs, title, bullets, subbullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x47, 0x5d)
    # Body
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.space_after = Pt(8)
        if subbullets and i < len(subbullets) and subbullets[i]:
            for sub in subbullets[i]:
                sp = tf.add_paragraph()
                sp.text = "  → " + sub
                sp.font.size = Pt(12)
                sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                sp.space_after = Pt(4)
    return slide

def add_bullet_slide(prs, title, bullets):
    return add_section_slide(prs, title, bullets)

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x47, 0x5d)
    # Simple table as text (pptx table API is verbose; we use textbox table layout)
    left, top = Inches(0.5), Inches(1.15)
    col_w = 2.8
    for c, h in enumerate(headers):
        box = slide.shapes.add_textbox(left + c * col_w * Inches(1), top, Inches(col_w), Inches(0.35))
        box.text_frame.paragraphs[0].text = h
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.size = Pt(11)
    for r, row in enumerate(rows):
        y_in = 1.15 + 0.38 + r * 0.36
        for c, cell in enumerate(row):
            box = slide.shapes.add_textbox(left + c * col_w * Inches(1), Inches(y_in), Inches(col_w), Inches(0.34))
            tf = box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = cell[:120] + ("…" if len(cell) > 120 else "")
            tf.paragraphs[0].font.size = Pt(10)
    return slide

def add_quote_slide(prs, title, quote, attribution=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x47, 0x5d)
    qbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(3))
    qbox.text_frame.word_wrap = True
    qbox.text_frame.paragraphs[0].text = quote
    qbox.text_frame.paragraphs[0].font.size = Pt(18)
    qbox.text_frame.paragraphs[0].font.italic = True
    if attribution:
        abox = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(8.6), Inches(0.5))
        abox.text_frame.paragraphs[0].text = attribution
        abox.text_frame.paragraphs[0].font.size = Pt(12)
        abox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    return slide

def add_closing_slide(prs, title, main_text, bullet_lines=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(28)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x47, 0x5d)
    bbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(2))
    bbox.text_frame.word_wrap = True
    bbox.text_frame.paragraphs[0].text = main_text
    bbox.text_frame.paragraphs[0].font.size = Pt(16)
    if bullet_lines:
        for line in bullet_lines:
            p = bbox.text_frame.add_paragraph()
            p.text = "• " + line
            p.font.size = Pt(14)
            p.space_before = Pt(6)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # MAX 9 SLIDES

    # 1. Title
    add_title_slide(
        prs,
        "System Dynamics Model Factory",
        "AeroDyn Systems — Strategic options for lethal AI & long-term business\nExternal System Modeling & AI Task Force"
    )

    # 2. Problem + Question (merged)
    add_bullet_slide(
        prs,
        "What the board asked for — and the question we focus on",
        [
            "Not one giant model; a repeatable way to build small, focused models for different strategic questions, with as little friction as possible.",
            "Focus: What happens to our long-term business if we heavily invest in lethal AI? How do public opinion, regulation, and contracts interact over 5–10 years?",
            "Mechanism: Lethal AI visibility → backlash → regulation → reputation → contract pipeline. Scenarios: base, high investment, stronger backlash, diversification.",
        ],
    )

    # 3. Reframing the question (bonus)
    add_section_slide(
        prs,
        "Reframing the question — from yes/no to conditions",
        [
            "Typical: “Should we invest in AI-controlled lethal weapons?” (yes/no) → doesn’t help explore when it pays off or backfires.",
            "Reframe: “Under what conditions does investing in lethal AI improve or harm our long-term position?”",
        ],
        [
            [],
            ["Mechanisms + levers (pace, diversification, transparency) + scenarios = strategic options, not a single number."],
        ],
    )

    # 4. What you get + How it works (merged)
    add_bullet_slide(
        prs,
        "What you get: insight, speed, reliability — and how it works",
        [
            "Insight: One question per model; see backlash, regulation, reputation, contracts over 5–10 years. Scenarios in minutes.",
            "Speed: Model Brief = contract (one page). New question → new model in hours to days. Same engine for all.",
            "Reliability: Explicit assumptions (stocks, flows, parameters); editable and auditable. No black box.",
            "Process: One question → one model. Building blocks reused. Scenarios, not point forecasts. Transparent and updatable.",
        ],
    )

    # 5. Key mechanisms
    add_bullet_slide(
        prs,
        "Key mechanisms the model captures",
        [
            "Lethal AI visibility → Public/ethical backlash (reacts to visibility; decays with media cycle).",
            "Backlash → Regulatory pressure (export controls, certification); eases with policy cycle.",
            "Backlash erodes Reputation (license to operate); transparency/compliance recover it.",
            "Reputation builds Contract pipeline; delivery drains it; regulation constrains it.",
        ],
    )

    # 6. Architecture + Systematic vs judgment (merged)
    add_bullet_slide(
        prs,
        "Architecture & what’s systematic vs human judgment",
        [
            "Repeatable: Single schema (stocks, flows, parameters); one engine builds ODEs and diagram from it. JSON in/out; GenAI can draft schema, human loads and validates.",
            "Systematic: Schema, simulation, diagram, import/export. Same inputs → same results.",
            "Human: What to include/exclude, how to interpret conflicting info, time horizon, parameter choices, reinterpretation of the question (e.g. yes/no → under what conditions).",
        ],
    )

    # 7. Trade-offs + What we don’t promise (merged)
    add_bullet_slide(
        prs,
        "Trade-offs & what we don’t promise",
        [
            "Simplicity over realism: small, explainable models; add another model for another question. GenAI assists; humans validate; provenance visible.",
            "No calibration on hard data yet: scenario-based exploration; assumptions documented.",
            "No perfect predictions; no single right answer. Models help you reason and compare options; choices (scope, interpretation, horizon) stay with you.",
        ],
    )

    # 8. Demo + Bottom line (merged)
    add_bullet_slide(
        prs,
        "Demo & bottom line",
        [
            "15-min demo: Pick question → load model (stocks, flows, loops) → change params or scenario → run → see “So what?” interpretation. Just the picture and the numbers.",
            "Bottom line: Credible, repeatable way to turn strategic questions into simple system-dynamics models the board can understand. Same approach for the next question — a library of small models, human choices explicit.",
        ],
    )

    # 9. Thank you
    add_title_slide(
        prs,
        "Thank you",
        "Questions — and we can run a live demo anytime.",
    )

    out_path = "AeroDyn_System_Dynamics_Model_Factory.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path} (9 slides)")

if __name__ == "__main__":
    main()
